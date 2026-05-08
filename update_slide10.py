"""Rebuild slide 10 (Tech Stack & Production Path) with a polished two-panel design."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

PPT = r'c:\Hackathon\VidyutDrishti\VidyutDrishti_Pitch.pptx'

# ── colour tokens ──────────────────────────────────────────────
C_BG        = RGBColor(248, 250, 252)   # near-white page bg
C_DARK      = RGBColor(15,  23,  42)    # slate-900
C_BODY      = RGBColor(30,  41,  59)    # slate-800
C_BLUE      = RGBColor(30,  64, 175)    # blue-800 (prototype panel)
C_BLUE_LT   = RGBColor(219, 234, 254)   # blue-100 row fill
C_BLUE_MED  = RGBColor(147, 197, 253)   # blue-300 accent
C_GREEN     = RGBColor(21, 128,  61)    # green-700 (production panel)
C_GREEN_LT  = RGBColor(220, 252, 231)   # green-100 row fill
C_GREEN_MED = RGBColor(134, 239, 172)   # green-300 accent
C_WHITE     = RGBColor(255, 255, 255)
C_GOLD      = RGBColor(161,  98,   7)   # amber-700
C_GOLD_BG   = RGBColor(254, 243, 199)   # amber-100
C_BORDER    = RGBColor(203, 213, 225)   # slate-300

PROTO_ROWS = [
    ('Language',    'Python 3.11  /  TypeScript'),
    ('API',         'FastAPI + Pydantic'),
    ('Database',    'TimescaleDB  —  single node, Docker'),
    ('Forecasting', 'Seasonal Baseline (STL + holidays + lags)'),
    ('Anomaly ML',  'Isolation Forest  —  scikit-learn'),
    ('Frontend',    'React 18 + Vite + Recharts + Leaflet'),
    ('Container',   'Docker Compose  —  3 services'),
    ('Auth',        'None  (prototype scope)'),
    ('Data',        'Synthetic AMI  —  259k readings, 8 DTs'),
]

PROD_ROWS = [
    ('Database',    'TimescaleDB distributed cluster (HA + replication)'),
    ('Ingestion',   'Kafka / Flink  —  replaces batch polling'),
    ('Auth',        'LDAP / RBAC  —  Section Eng / Inspector / Manager'),
    ('Deployment',  'Kubernetes  —  BESCOM data centre or NIC cloud'),
    ('Monitoring',  'Prometheus + Grafana + PagerDuty alerts'),
    ('Forecasting', 'Same model  —  retrained weekly on real AMI data'),
    ('ML Retrain',  'Isolation Forest  —  auto-retrained on inspector feedback'),
    ('AMI Link',    'Read-only REST/SFTP from real BESCOM AMI/MDM'),
    ('Scale',       '8.5M consumers  ·  50K DTs  ·  ~10M readings/day'),
]

# ── helpers ────────────────────────────────────────────────────
def rect(slide, l, t, w, h, fill, line_rgb=None, line_w=None):
    sh = slide.shapes.add_shape(1, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line_rgb:
        sh.line.color.rgb = line_rgb
        if line_w: sh.line.width = line_w
    else:
        sh.line.fill.background()
    return sh

def tb(slide, text, l, t, w, h, bold=False, size=10, color=C_BODY,
       align=PP_ALIGN.LEFT, wrap=True, italic=False):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame; tf.word_wrap = wrap
    p   = tf.paragraphs[0]; p.alignment = align
    r   = p.add_run(); r.text = text
    r.font.bold   = bold
    r.font.italic = italic
    r.font.size   = Pt(size)
    r.font.color.rgb = color
    return box

def draw_panel(slide, title, badge_txt, rows,
               panel_left, panel_top, panel_w, panel_h,
               hdr_rgb, hdr_lt, accent_rgb, badge_rgb, badge_txt_rgb):
    """Draw a labelled panel with a header badge and alternating rows."""
    ROW_H  = Inches(0.295)
    HDR_H  = Inches(0.50)
    KEY_W  = Inches(1.55)
    PAD    = Inches(0.08)

    # panel shadow (offset rect)
    rect(slide, panel_left + Inches(0.04), panel_top + Inches(0.04),
         panel_w, panel_h, RGBColor(203, 213, 225))

    # white panel card
    rect(slide, panel_left, panel_top, panel_w, panel_h,
         C_WHITE, C_BORDER, Pt(0.75))

    # header bar
    rect(slide, panel_left, panel_top, panel_w, HDR_H, hdr_rgb)

    # badge pill
    pill_w = Inches(1.5)
    pill_h = Inches(0.28)
    rect(slide, panel_left + PAD,
         panel_top + (HDR_H - pill_h) / 2, pill_w, pill_h, badge_rgb)
    tb(slide, badge_txt,
       panel_left + PAD + Inches(0.06),
       panel_top + (HDR_H - pill_h) / 2 + Inches(0.04),
       pill_w - Inches(0.12), pill_h - Inches(0.08),
       bold=True, size=7.5, color=badge_txt_rgb, align=PP_ALIGN.CENTER)

    # panel title
    tb(slide, title,
       panel_left + pill_w + PAD * 2,
       panel_top + Inches(0.06),
       panel_w - pill_w - PAD * 3, HDR_H - Inches(0.12),
       bold=True, size=12, color=C_WHITE, align=PP_ALIGN.LEFT)

    # rows
    for i, (key, val) in enumerate(rows):
        ry = panel_top + HDR_H + i * ROW_H
        bg = C_WHITE if i % 2 == 0 else hdr_lt

        # row background
        rect(slide, panel_left, ry, panel_w, ROW_H, bg)

        # left accent stripe
        rect(slide, panel_left, ry, Inches(0.04), ROW_H, accent_rgb)

        # key label (bold, dark)
        tb(slide, key,
           panel_left + Inches(0.10), ry + Inches(0.04),
           KEY_W, ROW_H - Inches(0.08),
           bold=True, size=8, color=hdr_rgb)

        # value
        tb(slide, val,
           panel_left + KEY_W + Inches(0.14), ry + Inches(0.04),
           panel_w - KEY_W - Inches(0.24), ROW_H - Inches(0.08),
           bold=False, size=8, color=C_BODY)

        # thin bottom hairline
        rect(slide, panel_left, ry + ROW_H - Pt(0.5),
             panel_w, Pt(0.5), C_BORDER)

# ── main ──────────────────────────────────────────────────────
prs   = Presentation(PPT)
W     = prs.slide_width
H     = prs.slide_height
slide = prs.slides[9]   # slide 10

# remove everything except the dark header bar shapes and footer/counter texts
KEEP_TEXTS = {
    'TECH STACK & PRODUCTION PATH',
    'Prototype uses real algorithms; production swaps infra — not logic',
    'VidyutDrishti  ·  AI for Bharat  ·  Theme 8: Smart Meter Intelligence & Loss Detection (BESCOM)',
    'github.com/srinidhirepala/VidyutDrishti',
}
sp_tree = slide.shapes._spTree
to_rm   = []
for shape in slide.shapes:
    t = shape.text_frame.text.strip() if shape.has_text_frame else ''
    if t in KEEP_TEXTS:
        continue
    to_rm.append(shape)

for shape in to_rm:
    sp_tree.remove(shape._element)

# ── light background fill for content area ──
CONTENT_TOP = Inches(1.15)
rect(slide, 0, CONTENT_TOP, W, H - CONTENT_TOP, C_BG)

# ── two panels ──
MARGIN   = Inches(0.28)
GAP      = Inches(0.22)
PANEL_W  = (W - MARGIN * 2 - GAP) / 2
PANEL_H  = Inches(4.15)
PANEL_T  = Inches(1.22)

draw_panel(
    slide, 'Prototype Stack', 'SUBMITTED',
    PROTO_ROWS,
    panel_left=MARGIN, panel_top=PANEL_T,
    panel_w=PANEL_W,   panel_h=PANEL_H,
    hdr_rgb=C_BLUE, hdr_lt=C_BLUE_LT, accent_rgb=C_BLUE_MED,
    badge_rgb=C_WHITE, badge_txt_rgb=C_BLUE,
)

draw_panel(
    slide, 'Production Changes', 'SAME LOGIC · BETTER INFRA',
    PROD_ROWS,
    panel_left=MARGIN + PANEL_W + GAP, panel_top=PANEL_T,
    panel_w=PANEL_W,                   panel_h=PANEL_H,
    hdr_rgb=C_GREEN, hdr_lt=C_GREEN_LT, accent_rgb=C_GREEN_MED,
    badge_rgb=C_WHITE, badge_txt_rgb=C_GREEN,
)

# ── bottom callout banner ──
CALL_T = PANEL_T + PANEL_H + Inches(0.18)
CALL_H = Inches(0.42)
rect(slide, MARGIN, CALL_T, W - MARGIN * 2, CALL_H, C_GOLD_BG, C_GOLD, Pt(1))
tb(slide,
   '✦  Detection & forecasting algorithms are IDENTICAL in prototype and production  —  only infrastructure and data sources change.',
   MARGIN + Inches(0.14), CALL_T + Inches(0.07),
   W - MARGIN * 2 - Inches(0.28), CALL_H - Inches(0.14),
   bold=True, size=9, color=C_GOLD, align=PP_ALIGN.CENTER)

prs.save(PPT)
print('Slide 10 rebuilt.')

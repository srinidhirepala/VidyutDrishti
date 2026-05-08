"""Rebuild slide 9 — clean magazine layout: large screenshots, thin caption, no clutter."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

SNAP = r'c:\Hackathon\VidyutDrishti\snapshots'
PPT  = r'c:\Hackathon\VidyutDrishti\VidyutDrishti_Pitch.pptx'

MODULES = [
    ('01-dashboard.png',         '01  Dashboard',        'KPIs · charts · forecast'),
    ('02-inspection-queue.png',  '02  Inspection Queue', 'Rs. × confidence ranking'),
    ('03-meter-lookup.png',      '03  Meter Lookup',     '4-layer drill-down'),
    ('04-feedback-form.png',     '04  Feedback',         'Live refresh on submit'),
    ('05-zone-map.png',          '05  Zone Risk Map',    'Bengaluru Leaflet heatmap'),
    ('06-evaluation-metrics.png','06  Eval Metrics',     'P=78% · R=85% · F1=0.81'),
    ('07-roi-calculator.png',    '07  ROI Calculator',   'Rs.455 Cr · 5-yr NPV'),
]

C_BG      = RGBColor(15,  23,  42)   # dark page bg same as header
C_CARD    = RGBColor(255, 255, 255)
C_CAPTION = RGBColor(15,  23,  42)   # near-black caption bar
C_TITLE   = RGBColor(255, 255, 255)
C_SUB     = RGBColor(147, 197, 253)  # blue-300
C_BORDER  = RGBColor(51,  65,  85)   # slate-700
C_ACCENT  = RGBColor(30,  64, 175)   # blue-800

prs   = Presentation(PPT)
W     = prs.slide_width
H     = prs.slide_height
slide = prs.slides[8]

# ── strip all old content except header texts ──
KEEP = {
    'WORKING PROTOTYPE',
    'Seven modules, all live in the submitted Docker stack',
    '09 / 12',
    'VidyutDrishti  ·  AI for Bharat  ·  Theme 8: Smart Meter Intelligence & Loss Detection (BESCOM)',
    'github.com/srinidhirepala/VidyutDrishti',
}
sp = slide.shapes._spTree
for shape in list(slide.shapes):
    t = shape.text_frame.text.strip() if shape.has_text_frame else ''
    if shape.shape_type == 13 or (t and t not in KEEP):
        sp.remove(shape._element)

# ── dark background fill below header ──
def rect(s, l, t, w, h, fill, lc=None, lw=None):
    sh = s.shapes.add_shape(1, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if lc: sh.line.color.rgb = lc; sh.line.width = lw or Pt(0.5)
    else:  sh.line.fill.background()
    return sh

def tb(s, text, l, t, w, h, bold=False, size=9, color=C_TITLE, align=PP_ALIGN.LEFT):
    box = s.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame; tf.word_wrap = False
    p   = tf.paragraphs[0]; p.alignment = align
    r   = p.add_run(); r.text = text
    r.font.bold = bold; r.font.size = Pt(size); r.font.color.rgb = color

CONTENT_TOP = Inches(1.18)
rect(slide, 0, CONTENT_TOP, W, H - CONTENT_TOP, C_BG)

# ── grid: 4 top + 3 bottom, maximise image area ──
PAD    = Inches(0.20)
GAP    = Inches(0.14)
CAP_H  = Inches(0.38)   # thin caption strip only

total_w  = W - PAD * 2
cw4      = (total_w - GAP * 3) / 4
cw3      = (total_w - GAP * 2) / 3

# two rows fit between CONTENT_TOP + PAD and bottom margin
usable_h = H - CONTENT_TOP - PAD * 2 - Inches(0.30)  # leave room for footer
row_h    = (usable_h - GAP) / 2
img_h    = row_h - CAP_H

row1_top = CONTENT_TOP + PAD
row2_top = row1_top + row_h + GAP

def add_card(img_path, title, subtitle, left, top, cw):
    # subtle shadow
    rect(slide, left + Inches(0.04), top + Inches(0.04), cw, row_h,
         RGBColor(8, 15, 30))
    # card white bg
    rect(slide, left, top, cw, row_h, C_CARD, C_BORDER, Pt(0.4))
    # screenshot
    slide.shapes.add_picture(img_path, left, top, cw, img_h)
    # caption bar
    rect(slide, left, top + img_h, cw, CAP_H, C_CAPTION)
    # blue left accent on caption
    rect(slide, left, top + img_h, Inches(0.05), CAP_H, C_ACCENT)
    # title in caption
    tb(slide, title,
       left + Inches(0.10), top + img_h + Inches(0.04),
       cw - Inches(0.12), Inches(0.18),
       bold=True, size=8.5, color=C_TITLE)
    # subtitle in caption
    tb(slide, subtitle,
       left + Inches(0.10), top + img_h + Inches(0.21),
       cw - Inches(0.12), Inches(0.15),
       bold=False, size=7, color=C_SUB)

for i in range(4):
    left = PAD + i * (cw4 + GAP)
    img, title, sub = MODULES[i]
    add_card(os.path.join(SNAP, img), title, sub, left, row1_top, cw4)

# centre the 3-card bottom row
row2_total = cw3 * 3 + GAP * 2
row2_left  = (W - row2_total) / 2
for i in range(3):
    left = row2_left + i * (cw3 + GAP)
    img, title, sub = MODULES[4 + i]
    add_card(os.path.join(SNAP, img), title, sub, left, row2_top, cw3)

prs.save(PPT)
print("Slide 9 rebuilt — clean layout.")
